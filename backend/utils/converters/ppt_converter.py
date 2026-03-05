from .base import BaseConverter
from pptx import Presentation


class PPTConverter(BaseConverter):
    """PowerPoint 文档转换器"""

    @classmethod
    def supports(cls, filename: str) -> bool:
        """判断是否支持该文件类型"""
        return filename.lower().endswith((".pptx", ".ppt"))

    @classmethod
    async def convert(cls, file_path: str) -> tuple[str, str]:
        """
        转换 PowerPoint 文档为 Markdown
        :return: (markdown_content, title)
        """
        prs = Presentation(file_path)
        markdown_lines = []
        title = "未命名演示文稿"

        for slide_idx, slide in enumerate(prs.slides):
            markdown_lines.append(f"\n## 幻灯片 {slide_idx + 1}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if slide_idx == 0 and not markdown_lines[2:]:
                        # 第一张幻灯片的第一个文本作为标题
                        title = text
                    markdown_lines.append(text)

        content = "\n\n".join(markdown_lines)
        return content, title
